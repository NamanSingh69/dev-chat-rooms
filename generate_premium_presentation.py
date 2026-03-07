import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors
DARK_BG = RGBColor(15, 15, 20)
PRIMARY_BLUE = RGBColor(59, 130, 246)
ACCENT_CYAN = RGBColor(34, 211, 238)
TEXT_WHITE = RGBColor(250, 250, 250)
TEXT_MUTED = RGBColor(160, 165, 180)

BACKGROUND_IMG = r"C:\Users\namsi\.gemini\antigravity\brain\08553432-e488-41da-be38-0105e8b6513e\low_key_dark_bg_1772226925421.png"
WORKFLOW_IMG = r"C:\Users\namsi\.gemini\antigravity\brain\08553432-e488-41da-be38-0105e8b6513e\wide_workflow_diagram_1772227240530.png"
UI_LOGIN = r"C:\Users\namsi\.gemini\antigravity\brain\08553432-e488-41da-be38-0105e8b6513e\.system_generated\click_feedback\click_feedback_1772213504774.png"
UI_CHAT = r"C:\Users\namsi\.gemini\antigravity\brain\08553432-e488-41da-be38-0105e8b6513e\.system_generated\click_feedback\click_feedback_1772213569790.png"

def add_bg(slide):
    if os.path.exists(BACKGROUND_IMG):
        slide.shapes.add_picture(BACKGROUND_IMG, 0, 0, width=Inches(16), height=Inches(9))
    else:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

def add_custom_text(slide, text, left, top, width, height, font_size, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    return p

def add_slide_title(slide, text):
    add_custom_text(slide, text, Inches(1), Inches(0.5), Inches(14), Inches(1), 44, TEXT_WHITE, True)

# Slide Layout 6 is Blank
BLANK_SLIDE = prs.slide_layouts[6]

# ----------------- SLIDE 1: Title -----------------
slide = prs.slides.add_slide(BLANK_SLIDE)
add_bg(slide)
add_custom_text(slide, "AI-Powered", Inches(1), Inches(3), Inches(14), Inches(1), 72, TEXT_WHITE, True, PP_ALIGN.CENTER)
add_custom_text(slide, "DEV CHAT ROOMS", Inches(1), Inches(4.2), Inches(14), Inches(1), 80, ACCENT_CYAN, True, PP_ALIGN.CENTER)
add_custom_text(slide, "Architecting Seamless Collaboration & LowEntropyAI", Inches(1), Inches(5.8), Inches(14), Inches(1), 28, TEXT_MUTED, False, PP_ALIGN.CENTER)

# ----------------- SLIDE 2: Value Proposition -----------------
slide = prs.slides.add_slide(BLANK_SLIDE)
add_bg(slide)
add_slide_title(slide, "The Problem & Value Proposition")

add_custom_text(slide, "The Gap in Developer Collaboration:", Inches(1), Inches(2.5), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "Context switching between code editors, chat apps, and AI tabs breaks flow.", Inches(1.5), Inches(3.2), Inches(13), Inches(0.5), 24, TEXT_WHITE)

add_custom_text(slide, "Our Value Proposition:", Inches(1), Inches(5.0), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "A unified workspace where context is king. Developers collaborate in real-time rooms.", Inches(1.5), Inches(5.7), Inches(13), Inches(0.5), 24, TEXT_WHITE)
add_custom_text(slide, "LowEntropyAI sits intimately within the chat, reading exact context to solve bugs instantly.", Inches(1.5), Inches(6.4), Inches(13), Inches(0.5), 24, TEXT_WHITE)


# ----------------- SLIDE 3: System Architecture -----------------
slide = prs.slides.add_slide(BLANK_SLIDE)
add_bg(slide)
add_slide_title(slide, "System Architecture & Workflow")
add_custom_text(slide, "Next.js Frontend connects via Socket.IO directly to the AI Core.", Inches(1), Inches(1.5), Inches(14), Inches(0.5), 24, TEXT_MUTED, False, PP_ALIGN.CENTER)

if os.path.exists(WORKFLOW_IMG):
    # Wide 16:9 layout taking up most of the slide
    slide.shapes.add_picture(WORKFLOW_IMG, Inches(1), Inches(2.2), width=Inches(14))


# ----------------- SLIDE 4: UX Principle -----------------
slide = prs.slides.add_slide(BLANK_SLIDE)
add_bg(slide)
add_slide_title(slide, "UX Principle: Glassmorphism")

add_custom_text(slide, "Why Glassmorphism?", Inches(1), Inches(2.5), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "Creates depth without clutter. By using blurred, translucent backgrounds (backdrop-blur), \nusers maintain spatial awareness without being distracted.", Inches(1.5), Inches(3.2), Inches(13), Inches(1), 24, TEXT_WHITE)

add_custom_text(slide, "Dark Mode Defaults:", Inches(1), Inches(5.0), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "Reduces eye strain for developers working long hours in low-light environments.", Inches(1.5), Inches(5.7), Inches(13), Inches(0.5), 24, TEXT_WHITE)

add_custom_text(slide, "High Contrast Actions:", Inches(1), Inches(7.0), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "Primary call-to-actions utilize strict #3b82f6 (Blue) ensuring AAA contrast ratings.", Inches(1.5), Inches(7.7), Inches(13), Inches(0.5), 24, TEXT_WHITE)


# ----------------- SLIDE 5: Frontend Showcase - Login -----------------
slide = prs.slides.add_slide(BLANK_SLIDE)
add_bg(slide)
add_slide_title(slide, "Frictionless Authentication UI")

# Left side text
add_custom_text(slide, "Modal-Driven Auth", Inches(1), Inches(3.0), Inches(5), Inches(0.5), 36, PRIMARY_BLUE, True)
add_custom_text(slide, "Captures Display Name & API Key gracefully without heavy page redirects.", Inches(1), Inches(4.0), Inches(4.5), Inches(2), 24, TEXT_WHITE)
add_custom_text(slide, "Clear error boundaries and link to Google AI Studio for key generation.", Inches(1), Inches(5.5), Inches(4.5), Inches(2), 24, TEXT_MUTED)

# Right side image
if os.path.exists(UI_LOGIN):
    # Constrain height to preserve aspect ratio and fit on slide
    slide.shapes.add_picture(UI_LOGIN, Inches(7.5), Inches(2.0), height=Inches(6.0))


# ----------------- SLIDE 6: Frontend Showcase - Chat -----------------
slide = prs.slides.add_slide(BLANK_SLIDE)
add_bg(slide)
add_slide_title(slide, "Core UI: The Collaboration Room")

# Right side text
add_custom_text(slide, "Real-Time Workspace", Inches(9.5), Inches(3.0), Inches(5.5), Inches(0.5), 36, PRIMARY_BLUE, True)
add_custom_text(slide, "Socket.IO chat interface featuring instantaneous broadcast to all clients.", Inches(9.5), Inches(4.0), Inches(5.5), Inches(2), 24, TEXT_WHITE)
add_custom_text(slide, "Markdown rendering and syntax highlighting built directly into the UI.", Inches(9.5), Inches(5.5), Inches(5.5), Inches(2), 24, TEXT_MUTED)

# Left side image
if os.path.exists(UI_CHAT):
    # Constrain height to preserve aspect ratio and fit on slide
    slide.shapes.add_picture(UI_CHAT, Inches(1), Inches(2.0), height=Inches(6.0))


# ----------------- SLIDE 7: The "A -> B -> C" User Flow -----------------
slide = prs.slides.add_slide(BLANK_SLIDE)
add_bg(slide)
add_slide_title(slide, "Streamlined User Flow (A -> B -> C)")

add_custom_text(slide, "[A] Secure Onboarding", Inches(1), Inches(2.5), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "User lands on app. Chat is visible but disabled. Click \"Kindly Login\" -> Modal opens.", Inches(1.5), Inches(3.2), Inches(13), Inches(0.5), 24, TEXT_WHITE)

add_custom_text(slide, "[B] Contextual Routing", Inches(1), Inches(4.5), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "User selects 'Frontend Issues' room. Socket.IO emits join, fetches historical messages.", Inches(1.5), Inches(5.2), Inches(13), Inches(0.5), 24, TEXT_WHITE)

add_custom_text(slide, "[C] AI Invocation", Inches(1), Inches(6.5), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "User tags @LowEntropyAI. AI reads entire room context and streams markdown response.", Inches(1.5), Inches(7.2), Inches(13), Inches(0.5), 24, TEXT_WHITE)


# ----------------- SLIDE 8: Technical Architecture -----------------
slide = prs.slides.add_slide(BLANK_SLIDE)
add_bg(slide)
add_slide_title(slide, "Technical Feasibility & Deployment")

add_custom_text(slide, "Tech Stack:", Inches(1), Inches(2.5), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "Next.js 15, Tailwind CSS V4, Socket.IO, @google/genai", Inches(1.5), Inches(3.2), Inches(13), Inches(0.5), 24, TEXT_WHITE)

add_custom_text(slide, "Dual-Mode Flexibility:", Inches(1), Inches(4.5), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "Local: Custom Node.js server handles WebSockets.", Inches(1.5), Inches(5.2), Inches(13), Inches(0.5), 24, TEXT_WHITE)
add_custom_text(slide, "Vercel: Falls back to Next.js serverless API routes for production stability.", Inches(1.5), Inches(5.9), Inches(13), Inches(0.5), 24, TEXT_MUTED)

add_custom_text(slide, "Cascading AI Fallback:", Inches(1), Inches(7.2), Inches(14), Inches(0.5), 32, PRIMARY_BLUE, True)
add_custom_text(slide, "Automatically shifts from gemini-3-flash -> gemini-2.5-pro -> flash if rate limits are hit.", Inches(1.5), Inches(7.9), Inches(13), Inches(0.5), 24, TEXT_WHITE)


# ----------------- SLIDE 9: Conclusion -----------------
slide = prs.slides.add_slide(BLANK_SLIDE)
add_bg(slide)
add_custom_text(slide, "THANK YOU", Inches(1), Inches(4), Inches(14), Inches(1), 80, ACCENT_CYAN, True, PP_ALIGN.CENTER)
add_custom_text(slide, "Ready for Deployment", Inches(1), Inches(5.5), Inches(14), Inches(1), 32, TEXT_MUTED, False, PP_ALIGN.CENTER)

prs.save("Hackathon_Premium_Design_Submission.pptx")
print("Presentation generated successfully at Hackathon_Premium_Design_Submission.pptx")
