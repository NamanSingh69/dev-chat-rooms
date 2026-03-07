from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Slide 1: Title & Value Proposition
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "AI-Powered Dev Chat Rooms"
subtitle.text = "Structural Design & User Experience Philosophy\n\nValue Proposition: Seamless developer collaboration augmented by context-aware AI."

# Slide 2: Landing Screen & Core UI
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide_2.shapes.title, slide_2.placeholders[1]
title.text = "Core Feature Screens"
tf = body.text_frame
tf.text = "1. Landing & Chat Interface:"
p = tf.add_paragraph()
p.text = "- Glassmorphic frosted sidebar for room navigation (reduces cognitive load)."
p.level = 1
p2 = tf.add_paragraph()
p2.text = "- Primary CTA positioned strategically within chat to gate unauthorized access."
p2.level = 1
p3 = tf.add_paragraph()
p3.text = "2. Settings & API Configuration:"
p4 = tf.add_paragraph()
p4.text = "- Clean, modal-based friction to capture Google Auth and AI Studio APIs."
p4.level = 1
p5 = tf.add_paragraph()
p5.text = "- Uses minimal color system for focus."
p5.level = 1

# Slide 3: User Flow (A -> B -> C)
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide_3.shapes.title, slide_3.placeholders[1]
title.text = "User Flow: The 'A -> B -> C' Pathway"
tf = body.text_frame
tf.text = "[A] Authentication:"
p = tf.add_paragraph()
p.text = "Guest -> Google Login / Anon Setup -> Chat Access."
p.level = 1

p2 = tf.add_paragraph()
p2.text = "[B] Room Selection:"
p2.level = 0
p3 = tf.add_paragraph()
p3.text = "User clicks ?room link or Sidebar Hash -> Socket.IO emits join event."
p3.level = 1

p4 = tf.add_paragraph()
p4.text = "[C] Collaboration & AI Resolution:"
p4.level = 0
p5 = tf.add_paragraph()
p5.text = "User tags @LowEntropyAI -> Context window passed to Gemini fallback -> Markdown streaming."
p5.level = 1

# Slide 4: Design System
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title, body = slide_4.shapes.title, slide_4.placeholders[1]
title.text = "Design System & Architecture"
tf = body.text_frame
tf.text = "Typography & Hierarchy:"
p = tf.add_paragraph()
p.text = "- Inter (sans-serif) for high legibility in dense chat logs."
p.level = 1
p2 = tf.add_paragraph()
p2.text = "Color Palette (Dark Mode First):"
p2.level = 0
p3 = tf.add_paragraph()
p3.text = "- Background: #0a0a0a (reduces eye strain)"
p3.level = 1
p4 = tf.add_paragraph()
p4.text = "- Primary Action: #3b82f6 (Blue) ensuring AAA contrast"
p4.level = 1
p5 = tf.add_paragraph()
p5.text = "Components:"
p5.level = 0
p6 = tf.add_paragraph()
p6.text = "- Cards: Border white/5 with backdrop-blur for depth."
p6.level = 1

prs.save("Hackathon_Design_Submission.pptx")
print("Presentation generated successfully at Hackathon_Design_Submission.pptx")
